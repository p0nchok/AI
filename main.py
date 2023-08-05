import openai
import requests
from pptx import Presentation
import os
import tkinter as tk

# Use the OpenAI API key
openai.api_key = "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"




def gui_main():
    def print_text():
        print("Test")


    window = tk.Tk()
    window.geometry("500x500")

    lbl_greeting = tk.Label(text="PPTX AI!")
    lbl_greeting.pack()

    lbl_topic = tk.Label(text="What is your presentation about?")
    lbl_topic.pack()

    txtbox_topic = tk.Text(window)
    txtbox_topic.pack()

    btn_submit = tk.Button(window, text="Submit", command=print_text())
    btn_submit.pack()

    window.mainloop()
def main():
    # GUI Builder


    # AI questioning
    topic = input("What's the topic of your presentation? ")
    prs = Presentation()
    layout = prs.slide_layouts[1]

    stop = True
    topic_ai_question = f"Write key topics in a presentation about {topic}"
    topic_ai_answer = new_chatbot(topic_ai_question)

    sub_topics_list = topic_ai_answer.split('\n')
    sub_topics_list.pop(0)
    sub_topics_answers = []
    for sub_topic in sub_topics_list:
        sub_topics_answers.append(new_chatbot(sub_topic))

    for i in range(len(sub_topics_list)):
        slide = prs.slides.add_slide(layout)
        slide_title = slide.shapes.title
        slide_subtitle = slide.placeholders[1]

        title = sub_topics_list[i]
        title = title.replace("-", "")
        slide_title.text = title
        subtitle = sub_topics_answers[i]
        subtitle = subtitle.replace("-", "")
        subtitle = subtitle.replace("\n\n", "\n")
        slide_subtitle.text = subtitle

    # Main Page End

    # Slides:
    print(sub_topics_list)

    prs.save(fr'C:\Users\User\PycharmProjects\AI\Presentations\{topic}.pptx')


def new_chatbot(input):
    # Get user input
    user_input = input

    # Use the OpenAI API to generate a response
    response = openai.Completion.create(
        engine="text-davinci-002",
        prompt=(f"{user_input} \n"),
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.7
    )
    ai_output = response.choices[0].text
    return ai_output


# def imagebot():
#     desc = "How to use body language to persuade"
#     response = openai.Image.create(
#         prompt=desc,
#         n=1,
#         size="1024x1024"
#     )
#     image_url = response['data'][0]['url']
#
#     name = desc.replace(' ', '_')
#     name = f"{name}.jpg"
#     save_path = name
#     response = requests.get(image_url)
#     with open(save_path, "wb") as file:
#         # Write the contents of the response (r.content) to a new file
#         file.write(response.content)


# def powerpoint(title_ai, content):
#     print("PowerPoint Function")
#     prs = Presentation()
#     title_slide_layout = prs.slide_layouts[0]
#     slide = prs.slides.add_slide(title_slide_layout)
#     title = slide.shapes.title
#     subtitle = slide.placeholders[1]
#
#     # title.text = "Hello, World!"
#     title.text = title_ai
#     # subtitle.text = "python-pptx was here!"
#     subtitle.text = content
#
#     prs.save('test.pptx')

def image_generator():
    PROMPT = "An eco-friendly computer from the 90s in the style of vaporwave"

    openai.api_key = os.getenv("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")

    response = openai.Image.create(
        prompt=PROMPT,
        n=1,
        size="256x256",
    )

    print(response["data"][0]["url"])
if __name__ == "__main__":
    # main()
    gui_main()

    # chatbot()
    # imagebot()
    # powerpoint()
    # image_generator()
